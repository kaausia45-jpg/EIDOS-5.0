"""[2026-06-03 P4] AIRA 발표기 — 슬라이드 직접 생성 + Q&A(자율 back-nav) 핵심 로직.

순수 데이터/생성 로직(EIDOS GUI 무의존·LLM 주입형·테스트 가능). GUI(AiraDeckDialog·
레이저 포인터·발표 진행)는 eidos_chat_gui.py 에 둔다. python-pptx 내보내기 선택 제공.

설계:
- Deck = 제목 + Slide 목록. Slide = 제목 + 불릿 + AIRA 내레이션 + 표정 힌트.
- generate_deck_async: 주제 → LLM → 슬라이드 구조 JSON. 실패 시 안전한 1장 폴백.
- answer_question_async: Q&A. 답변 + 어느 슬라이드로 돌아가 설명할지(slide_index) 판단.
  → GUI 가 그 슬라이드로 자율 이동 후 레이저로 짚으며 설명(요구사항 #4).
- deck_to_pptx: python-pptx 로 .pptx 파일 내보내기(있을 때만).
"""
from __future__ import annotations

import json
import re
from dataclasses import dataclass, field

# 표정 키 힌트 — assets/expr_*.png 에 실제 존재하는 것들(생성기가 이 안에서 고르게 유도).
# GUI 가 실제 스캔 목록으로 최종 검증하므로 여기 누락돼도 안전.
_EXPR_HINTS = (
    "information", "point", "teaching", "excited", "happy",
    "confidence", "question", "thinking", "neutral", "ending", "surprised",
)


@dataclass
class Slide:
    title: str = ""
    bullets: list = field(default_factory=list)   # list[str]
    narration: str = ""                            # AIRA 가 말할 대사
    expr: str = "information"                       # 표정 키 힌트

    def to_dict(self) -> dict:
        return {
            "title": self.title,
            "bullets": list(self.bullets),
            "narration": self.narration,
            "expr": self.expr,
        }

    @staticmethod
    def from_dict(d: dict) -> "Slide":
        d = d or {}
        b = d.get("bullets") or []
        if isinstance(b, str):
            b = [b]
        return Slide(
            title=str(d.get("title", "") or "").strip(),
            bullets=[str(x).strip() for x in b if str(x).strip()][:6],
            narration=str(d.get("narration", "") or "").strip(),
            expr=str(d.get("expr", "information") or "information").strip().lower(),
        )


@dataclass
class Deck:
    title: str = ""
    slides: list = field(default_factory=list)     # list[Slide]
    topic: str = ""
    generated: bool = True                          # False = LLM 생성 실패 폴백(띄우지 말 것)

    def to_dict(self) -> dict:
        return {
            "title": self.title,
            "topic": self.topic,
            "slides": [s.to_dict() for s in self.slides],
            "generated": self.generated,
        }

    @staticmethod
    def from_dict(d: dict) -> "Deck":
        d = d or {}
        return Deck(
            title=str(d.get("title", "") or "").strip(),
            topic=str(d.get("topic", "") or "").strip(),
            slides=[Slide.from_dict(s) for s in (d.get("slides") or [])],
            generated=bool(d.get("generated", True)),
        )


# ── JSON 파싱(실 LLM 강건성) ────────────────────────────────────────────────
def _strip_fences(text: str) -> str:
    t = (text or "").strip()
    if t.startswith("```"):
        t = re.sub(r"^```[a-zA-Z0-9]*\s*", "", t)
        if t.endswith("```"):
            t = t[:-3]
    return t.strip()


def _extract_balanced(text: str, open_ch: str, close_ch: str) -> str:
    """문자열 리터럴 안의 괄호는 무시하고 균형 잡힌 첫 블록 추출."""
    s = text
    start = s.find(open_ch)
    if start < 0:
        return ""
    depth = 0
    in_str = False
    esc = False
    quote = ""
    for i in range(start, len(s)):
        c = s[i]
        if in_str:
            if esc:
                esc = False
            elif c == "\\":
                esc = True
            elif c == quote:
                in_str = False
            continue
        if c in ('"', "'"):
            in_str = True
            quote = c
        elif c == open_ch:
            depth += 1
        elif c == close_ch:
            depth -= 1
            if depth == 0:
                return s[start:i + 1]
    return ""


def _safe_json(text: str):
    t = _strip_fences(text)
    try:
        return json.loads(t)
    except Exception:
        pass
    block = _extract_balanced(t, "{", "}")
    if block:
        try:
            return json.loads(block)
        except Exception:
            pass
    # 잘린 JSON(토큰 한도 초과 등) — 완성된 슬라이드만 살려낸다.
    return _salvage_slides(t)


def _salvage_slides(text: str):
    """truncate 된 응답에서 완성된 슬라이드 객체만 추출(JSON 끝이 잘려도 가능한 만큼)."""
    t = _strip_fences(text or "")
    if '"slides"' not in t:
        return None
    mt = re.search(r'"title"\s*:\s*"([^"]*)"', t)
    title = mt.group(1) if mt else ""
    arr_start = t.find("[", t.find('"slides"'))
    if arr_start < 0:
        return None
    rest = t[arr_start + 1:]
    slides = []
    while True:
        obj = _extract_balanced(rest, "{", "}")
        if not obj:
            break
        try:
            slides.append(json.loads(obj))
        except Exception:
            pass
        nxt = rest.find(obj) + len(obj)
        rest = rest[nxt:]
        if "{" not in rest:
            break
    if slides:
        return {"title": title, "slides": slides}
    return None


# ── 슬라이드 생성 ───────────────────────────────────────────────────────────
async def generate_deck_async(
    topic: str,
    *,
    llm_async,
    n_slides: int = 6,
    audience: str = "",
    timeout_sec: float = 60.0,
) -> Deck:
    """주제 → AIRA 발표 슬라이드 덱(LLM). 실패 시 1장짜리 안전 폴백.

    llm_async: get_llm_response_async 호환(prompt, max_tokens=..., use_cache=...).
    """
    topic = (topic or "").strip()
    if not topic:
        return Deck(title="발표", slides=[
            Slide(title="발표 주제가 비어 있어요",
                  bullets=["발표할 주제를 알려주세요."],
                  narration="발표할 주제를 알려주시면 제가 슬라이드를 만들어 발표할게요.",
                  expr="question")], topic=topic)

    n = max(3, min(int(n_slides or 6), 12))
    aud = f"\n청중: {audience.strip()}" if (audience or "").strip() else ""
    prompt = (
        "너는 ENFP 발표자 AIRA 야. 아래 주제로 발표 슬라이드 덱을 JSON 으로 설계해.\n"
        f"주제: {topic}{aud}\n"
        f"- 슬라이드 {n}장 내외(표지 포함). 첫 장은 표지(주제·한줄소개).\n"
        "- 각 슬라이드: title(짧게), bullets(2~5개·핵심만·각 40자 이내), "
        "narration(AIRA 가 그 슬라이드에서 말할 대사·친근한 존댓말 해요체·반말 금지·2~4문장), "
        "expr(표정 키 1개).\n"
        f"- expr 는 다음 중 하나: {', '.join(_EXPR_HINTS)}.\n"
        "- 마지막 내용 슬라이드 뒤 Q&A 슬라이드는 추가하지 마(시스템이 자동으로 붙임).\n"
        "JSON 형식만 출력(설명 금지):\n"
        '{"title":"덱 제목","slides":[{"title":"...","bullets":["..."],'
        '"narration":"...","expr":"information"}]}'
    )
    import asyncio
    data = None
    raw = ""
    # 최대 2회 시도 — 토큰 한도 상향(4096)으로 JSON 잘림 방지 + 일시 실패 재시도.
    for attempt in range(2):
        try:
            raw = await asyncio.wait_for(
                llm_async(prompt, max_tokens=4096, use_cache=False), timeout=timeout_sec)
        except Exception as e:
            raw = ""
            print(f"[aira deck] generate 시도{attempt + 1} 실패: {e}")
        data = _safe_json(raw) if raw else None
        if isinstance(data, dict) and data.get("slides"):
            break
        print(f"[aira deck] 시도{attempt + 1}: 슬라이드 파싱 실패 → 재시도")

    if not isinstance(data, dict) or not data.get("slides"):
        # 생성 실패 — generated=False 로 표시(GUI 가 PowerPoint 로 띄우지 않고 안내).
        return Deck(
            title=topic[:60],
            topic=topic,
            generated=False,
            slides=[Slide(
                title=topic[:60],
                bullets=["슬라이드 자동 생성에 실패했어요.", "주제를 더 구체적으로 알려주시면 다시 만들어 볼게요."],
                narration=f"{topic} 에 대한 발표를 준비하려 했는데, 슬라이드 생성이 잘 안 됐어요. 주제를 조금 더 구체적으로 알려주시겠어요?",
                expr="concerned")],
        )

    deck = Deck.from_dict({**data, "topic": topic})
    if not deck.title:
        deck.title = topic[:60]
    # 빈 슬라이드 정리
    deck.slides = [s for s in deck.slides if (s.title or s.bullets or s.narration)]
    if not deck.slides:
        deck.slides = [Slide(title=topic[:60], bullets=[topic[:80]],
                             narration=f"{topic} 에 대해 발표할게요.", expr="information")]
    return deck


# ── Q&A — 답변 + 자율 back-nav 판단 ─────────────────────────────────────────
def deck_outline(deck: Deck, max_chars: int = 1800) -> str:
    """Q&A 프롬프트용 — 슬라이드 인덱스+제목+불릿 요약."""
    lines = []
    for i, s in enumerate(deck.slides):
        b = " / ".join(s.bullets[:4])
        lines.append(f"[{i}] {s.title} — {b}"[:200])
    out = "\n".join(lines)
    return out[:max_chars]


async def answer_question_async(
    deck: Deck,
    question: str,
    *,
    llm_async,
    timeout_sec: float = 40.0,
) -> dict:
    """발표 후 Q&A. 답변 + 어느 슬라이드로 돌아가 설명할지 판단.

    반환: {"answer": str, "slide_index": int(-1 이면 이동 안 함), "expr": str}.
    """
    question = (question or "").strip()
    if not question:
        return {"answer": "", "slide_index": -1, "expr": "question"}
    outline = deck_outline(deck)
    prompt = (
        "너는 방금 발표를 마친 ENFP 발표자 AIRA 야. 청중의 질문에 답해.\n"
        "친근한 존댓말 해요체로, 2~4문장 간결하게. 반말 금지.\n"
        "필요하면 발표 슬라이드 중 하나로 돌아가 짚으며 설명할 수 있어 — 그 슬라이드 번호를 "
        "slide_index 로 지정해(돌아갈 필요 없으면 -1).\n"
        f"[발표 슬라이드 목록]\n{outline}\n"
        f"[질문] {question}\n"
        "JSON 만 출력(설명 금지):\n"
        '{"answer":"...","slide_index":-1,"expr":"teaching"}'
    )
    raw = ""
    try:
        import asyncio
        raw = await asyncio.wait_for(
            llm_async(prompt, max_tokens=600, use_cache=False), timeout=timeout_sec)
    except Exception as e:
        print(f"[aira deck] qa 실패: {e}")
        raw = ""

    data = _safe_json(raw) if raw else None
    if not isinstance(data, dict):
        return {"answer": "죄송해요, 지금은 답변을 정리하지 못했어요. 다시 한 번 여쭤봐 주시겠어요?",
                "slide_index": -1, "expr": "concerned"}
    try:
        si = int(data.get("slide_index", -1))
    except Exception:
        si = -1
    if si < 0 or si >= len(deck.slides):
        si = -1
    return {
        "answer": str(data.get("answer", "") or "").strip(),
        "slide_index": si,
        "expr": str(data.get("expr", "teaching") or "teaching").strip().lower(),
    }


# ── PPTX 내보내기(선택·python-pptx 있을 때만) ───────────────────────────────
def pptx_available() -> bool:
    try:
        import pptx  # noqa: F401
        return True
    except Exception:
        return False


def deck_to_pptx(deck: Deck, path: str) -> str:
    """Deck → .pptx 파일. python-pptx 필요. 성공 시 path 반환, 실패 시 예외."""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    # 표지
    title_layout = prs.slide_layouts[0]
    s0 = prs.slides.add_slide(title_layout)
    s0.shapes.title.text = deck.title or (deck.topic or "발표")[:60]
    try:
        if len(s0.placeholders) > 1:
            s0.placeholders[1].text = "AIRA 발표"
    except Exception:
        pass

    bullet_layout = prs.slide_layouts[1]
    for sl in deck.slides:
        s = prs.slides.add_slide(bullet_layout)
        try:
            s.shapes.title.text = sl.title or " "
        except Exception:
            pass
        try:
            body = s.placeholders[1].text_frame
            body.clear()
            for i, b in enumerate(sl.bullets or [sl.narration or ""]):
                p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                p.text = b
                p.font.size = Pt(18)
            # 발표 노트에 내레이션 저장
            if sl.narration:
                s.notes_slide.notes_text_frame.text = sl.narration
        except Exception:
            pass

    # Q&A 표지
    try:
        q = prs.slides.add_slide(prs.slide_layouts[1])
        q.shapes.title.text = "Q&A"
        q.placeholders[1].text_frame.text = "궁금한 점을 물어보세요!"
    except Exception:
        pass

    prs.save(path)
    return path
